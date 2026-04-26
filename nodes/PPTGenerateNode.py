import asyncio
import json
import re
import os
from datetime import datetime
from typing import Dict, List

from langchain_core.messages import HumanMessage, SystemMessage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from state.state import IMState
from utils.logger_handler import logger
from nodes.agent.ppt_agent import ppt_outline_agent, ppt_outline_revision_agent, ppt_content_agent
from nodes.agent.prompt.ppt_prompt import ppt_outline_prompt, ppt_content_prompt, ppt_outline_revision_prompt
from utils.feishuUtils import feishu_api


def extract_json(content: str) -> str:
    """从 LLM 返回内容中提取 JSON"""
    json_match = re.search(r'```(?:json)?\s*([\s\S]*?)\s*```', content)
    if json_match:
        return json_match.group(1).strip()
    return content.strip()


def get_ppt_style_from_intent(intent: Dict) -> str:
    """根据用户意图和主题推断PPT风格"""
    topic = intent.get("topic", "").lower()
    
    tech_keywords = ["技术", "tech", "代码", "code", "算法", "ai", "人工智能", "开发", "架构", "系统"]
    creative_keywords = ["创意", "营销", "活动", "推广", "品牌", "广告", "设计"]
    edu_keywords = ["教育", "培训", "学习", "课程", "环保", "健康", "医疗"]
    luxury_keywords = ["高端", "时尚", "艺术", "设计", "品牌", "精品", "奢侈"]
    academic_keywords = ["学术", "论文", "研究", "报告", "分析", "极简", "简约"]
    
    if any(kw in topic for kw in tech_keywords):
        return "tech_dark"
    elif any(kw in topic for kw in creative_keywords):
        return "warm_orange"
    elif any(kw in topic for kw in edu_keywords):
        return "fresh_green"
    elif any(kw in topic for kw in luxury_keywords):
        return "elegant_purple"
    elif any(kw in topic for kw in academic_keywords):
        return "minimal_white"
    else:
        return "business_blue"


def build_ppt_outline_messages(state: IMState) -> list:
    """构建PPT大纲生成的消息列表"""
    intent = state.get("intent", {})
    chat_context = state.get("chat_context", "")
    
    topic = intent.get("topic", "未命名PPT")
    ppt_type = intent.get("additional_info", {}).get("ppt_type", "meeting_presentation")
    key_points = intent.get("key_points", [])
    ppt_style = get_ppt_style_from_intent(intent)
    
    current_date = datetime.now().strftime("%Y.%m.%d")
    key_points_str = "\n".join([f"- {point}" for point in key_points]) if key_points else "无特定要点"
    
    user_message = f"""请根据以下信息生成PPT大纲：

**PPT主题**: {topic}
**PPT类型**: {ppt_type}
**PPT风格**: {ppt_style}
**关键要点**:
{key_points_str}
**上下文信息**: {chat_context if chat_context else "无"}
**当前日期**: {current_date}

请生成结构清晰的PPT大纲。"""

    return [
        SystemMessage(content=ppt_outline_prompt),
        HumanMessage(content=user_message)
    ]


def build_ppt_outline_revision_messages(state: IMState) -> list:
    """构建PPT大纲修改的消息列表"""
    intent = state.get("intent", {})
    chat_context = state.get("chat_context", "")
    ppt_structure = state.get("ppt_structure", {})
    ppt_outline_feedback = state.get("ppt_outline_feedback", "")
    
    topic = intent.get("topic", "未命名PPT")
    ppt_type = ppt_structure.get("ppt_type", "meeting_presentation")
    ppt_style = ppt_structure.get("ppt_style", "business_blue")
    key_points = intent.get("key_points", [])
    
    current_date = datetime.now().strftime("%Y.%m.%d")
    key_points_str = "\n".join([f"- {point}" for point in key_points]) if key_points else "无特定要点"
    current_outline_str = json.dumps(ppt_structure, ensure_ascii=False, indent=2)
    
    user_message = f"""请根据用户反馈修改PPT大纲：

**PPT主题**: {topic}
**PPT类型**: {ppt_type}
**PPT风格**: {ppt_style}
**当前大纲**:
{current_outline_str}

**用户反馈**: {ppt_outline_feedback}

**关键要点**:
{key_points_str}
**上下文信息**: {chat_context if chat_context else "无"}
**当前日期**: {current_date}

请根据用户反馈修改PPT大纲。"""

    return [
        SystemMessage(content=ppt_outline_revision_prompt),
        HumanMessage(content=user_message)
    ]


def build_ppt_content_messages(ppt_structure: Dict, slide: Dict, chat_context: str) -> list:
    """构建PPT内容生成的消息列表"""
    title = ppt_structure.get("title", "")
    ppt_type = ppt_structure.get("ppt_type", "meeting_presentation")
    ppt_style = ppt_structure.get("ppt_style", "business_blue")
    
    page = slide.get("page", 1)
    slide_type = slide.get("type", "content")
    slide_title = slide.get("title", "")
    original_points = slide.get("points", [])
    
    points_str = "\n".join([f"- {point}" for point in original_points]) if original_points else "无"
    
    user_message = f"""请为以下PPT页面生成丰富内容：

**PPT主题**: {title}
**PPT类型**: {ppt_type}
**PPT风格**: {ppt_style}
**当前页面**: 第{page}页
**页面类型**: {slide_type}
**页面标题**: {slide_title}
**原始要点**:
{points_str}
**上下文信息**: {chat_context if chat_context else "无"}

请生成丰富、有数据支撑、适合视觉呈现的PPT内容。"""

    return [
        SystemMessage(content=ppt_content_prompt),
        HumanMessage(content=user_message)
    ]


def format_ppt_outline(ppt_structure: Dict) -> str:
    """格式化PPT大纲为易读的文本"""
    lines = []
    lines.append("=" * 50)
    lines.append(f"[PPT标题] {ppt_structure.get('title', 'N/A')}")
    lines.append(f"[PPT类型] {ppt_structure.get('ppt_type', 'N/A')}")
    lines.append(f"[PPT风格] {ppt_structure.get('ppt_style', 'business_blue')}")
    lines.append("=" * 50)
    
    lines.append("\n[PPT结构]")
    slides = ppt_structure.get('slides', [])
    for slide in slides:
        page = slide.get('page', 0)
        slide_type = slide.get('type', 'content')
        title = slide.get('title', '未命名页面')
        
        type_emoji = {
            'cover': '📘',
            'agenda': '📋',
            'content': '📝',
            'summary': '✅',
            'ending': '🎉'
        }.get(slide_type, '📄')
        
        lines.append(f"  {type_emoji} 第{page}页 [{slide_type}] {title}")
        
        if slide_type in ['agenda', 'content', 'summary']:
            points = slide.get('points', [])
            for point in points:
                lines.append(f"     • {point}")
        
        if slide.get('subtitle'):
            lines.append(f"     副标题: {slide['subtitle']}")
    
    lines.append("\n" + "=" * 50)
    lines.append("请确认此大纲是否满意？")
    lines.append("- 输入 '确认' 或 'ok' 开始生成PPT")
    lines.append("- 输入修改意见，如'增加一页关于xxx的内容'")
    lines.append("- 输入 '取消' 终止PPT生成")
    lines.append("=" * 50)
    
    return "\n".join(lines)


# 模板风格配色方案
PPT_STYLES = {
    "business_blue": {
        "name": "商务蓝",
        "primary": RGBColor(41, 98, 255),
        "secondary": RGBColor(0, 150, 199),
        "accent": RGBColor(255, 107, 107),
        "dark": RGBColor(33, 37, 41),
        "light": RGBColor(248, 249, 250),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(41, 98, 255),
    },
    "tech_dark": {
        "name": "科技黑",
        "primary": RGBColor(0, 255, 136),
        "secondary": RGBColor(0, 184, 255),
        "accent": RGBColor(255, 50, 100),
        "dark": RGBColor(18, 18, 18),
        "light": RGBColor(30, 30, 30),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(18, 18, 18),
    },
    "fresh_green": {
        "name": "清新绿",
        "primary": RGBColor(46, 204, 113),
        "secondary": RGBColor(39, 174, 96),
        "accent": RGBColor(241, 196, 15),
        "dark": RGBColor(44, 62, 80),
        "light": RGBColor(236, 240, 241),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(46, 204, 113),
    },
    "warm_orange": {
        "name": "活力橙",
        "primary": RGBColor(255, 140, 0),
        "secondary": RGBColor(255, 87, 34),
        "accent": RGBColor(255, 193, 7),
        "dark": RGBColor(62, 39, 35),
        "light": RGBColor(255, 243, 224),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(255, 140, 0),
    },
    "minimal_white": {
        "name": "极简白",
        "primary": RGBColor(52, 73, 94),
        "secondary": RGBColor(149, 165, 166),
        "accent": RGBColor(231, 76, 60),
        "dark": RGBColor(44, 62, 80),
        "light": RGBColor(250, 250, 250),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(52, 73, 94),
    },
    "elegant_purple": {
        "name": "优雅紫",
        "primary": RGBColor(155, 89, 182),
        "secondary": RGBColor(142, 68, 173),
        "accent": RGBColor(236, 112, 99),
        "dark": RGBColor(44, 62, 80),
        "light": RGBColor(245, 238, 250),
        "white": RGBColor(255, 255, 255),
        "background": RGBColor(155, 89, 182),
    },
}


def add_number_card(slide, x, y, width, height, number, label, color, bg_color):
    """添加数字卡片"""
    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()
    
    # 数字
    num_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(width), Inches(height * 0.5))
    num_frame = num_box.text_frame
    num_frame.text = str(number)
    num_para = num_frame.paragraphs[0]
    num_para.font.size = Pt(36)
    num_para.font.bold = True
    num_para.font.color.rgb = color
    num_para.alignment = PP_ALIGN.CENTER
    
    # 标签
    label_box = slide.shapes.add_textbox(Inches(x), Inches(y + height * 0.5), Inches(width), Inches(height * 0.4))
    label_frame = label_box.text_frame
    label_frame.text = label
    label_para = label_frame.paragraphs[0]
    label_para.font.size = Pt(14)
    label_para.font.color.rgb = RGBColor(108, 117, 125)
    label_para.alignment = PP_ALIGN.CENTER


def add_chart(slide, x, y, width, height, chart_type, data, colors):
    """添加图表"""
    try:
        chart_data = ChartData()
        chart_data.categories = data.get('labels', [])
        chart_data.add_series('数据', data.get('values', []))
        
        if chart_type == 'bar':
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), 
                Inches(width), Inches(height), chart_data
            ).chart
        elif chart_type == 'pie':
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(x), Inches(y), 
                Inches(width), Inches(height), chart_data
            ).chart
        elif chart_type == 'line':
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, Inches(x), Inches(y), 
                Inches(width), Inches(height), chart_data
            ).chart
        
        return chart
    except Exception as e:
        logger.warning(f"添加图表失败: {e}")
        return None


def create_ppt_file(ppt_structure: Dict, ppt_content: Dict, output_path: str) -> str:
    """使用python-pptx创建美观的PPT文件"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    style_name = ppt_structure.get("ppt_style", "business_blue")
    COLORS = PPT_STYLES.get(style_name, PPT_STYLES["business_blue"])
    
    slides_data = ppt_content.get("slides", [])
    
    for idx, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        content = slide_data.get("content", {})
        visual_suggestions = content.get("visual_suggestions", {})
        main_points = content.get("main_points", [])
        
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        if slide_type == "cover":
            # 封面页 - 渐变背景效果
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = COLORS['background']
            bg_shape.line.fill.background()
            
            # 装饰性圆形
            circle1 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(6), Inches(6)
            )
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = COLORS['white']
            circle1.line.fill.background()
            
            circle2 = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(-2), Inches(5), Inches(4), Inches(4)
            )
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = COLORS['secondary']
            circle2.line.fill.background()
            
            # 主标题
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['white']
            title_para.alignment = PP_ALIGN.CENTER
            
            # 副标题
            if subtitle:
                subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.733), Inches(1))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = subtitle
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(28)
                subtitle_para.font.color.rgb = COLORS['white']
                subtitle_para.alignment = PP_ALIGN.CENTER
            
            # 演讲者信息
            speaker_info = content.get("speaker_info", "")
            if speaker_info:
                info_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.733), Inches(0.8))
                info_frame = info_box.text_frame
                info_frame.text = speaker_info
                info_para = info_frame.paragraphs[0]
                info_para.font.size = Pt(18)
                info_para.font.color.rgb = COLORS['white']
                info_para.alignment = PP_ALIGN.CENTER
                
        elif slide_type == "agenda":
            # 目录页 - 左侧色块设计
            left_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5)
            )
            left_bar.fill.solid()
            left_bar.fill.fore_color.rgb = COLORS['primary']
            left_bar.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(12), Inches(1))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(42)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['dark']
            
            # 副标题
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(12), Inches(0.6))
                sub_frame = sub_box.text_frame
                sub_frame.text = subtitle
                sub_para = sub_frame.paragraphs[0]
                sub_para.font.size = Pt(18)
                sub_para.font.color.rgb = COLORS['gray'] if 'gray' in COLORS else RGBColor(108, 117, 125)
            
            # 目录项 - 双列布局
            points = content.get("points", [])
            if points:
                items_per_col = (len(points) + 1) // 2
                for i, point in enumerate(points[:8]):  # 最多8个
                    col = i // items_per_col
                    row = i % items_per_col
                    x = 0.8 + col * 6
                    y = 2.2 + row * 1.0
                    
                    # 序号
                    num_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.5), Inches(0.6))
                    num_frame = num_box.text_frame
                    num_frame.text = f"{i+1:02d}"
                    num_para = num_frame.paragraphs[0]
                    num_para.font.size = Pt(24)
                    num_para.font.bold = True
                    num_para.font.color.rgb = COLORS['primary']
                    
                    # 内容
                    content_box = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y), Inches(5), Inches(0.8))
                    content_frame = content_box.text_frame
                    content_frame.word_wrap = True
                    content_frame.text = point
                    content_para = content_frame.paragraphs[0]
                    content_para.font.size = Pt(20)
                    content_para.font.color.rgb = COLORS['dark']
                    
        elif slide_type == "content":
            # 内容页 - 多种布局支持
            layout_type = visual_suggestions.get("layout", "standard")
            
            # 顶部装饰条
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = COLORS['primary']
            top_bar.line.fill.background()
            
            # 标题区域
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['dark']
            
            # 副标题
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12), Inches(0.5))
                sub_frame = sub_box.text_frame
                sub_frame.text = subtitle
                sub_para = sub_frame.paragraphs[0]
                sub_para.font.size = Pt(16)
                sub_para.font.color.rgb = COLORS['gray'] if 'gray' in COLORS else RGBColor(108, 117, 125)
            
            # 根据布局类型渲染内容
            if layout_type == "three_columns" and len(main_points) >= 3:
                # 三栏布局
                col_width = 4
                for i, point in enumerate(main_points[:3]):
                    x = 0.6 + i * (col_width + 0.3)
                    
                    # 图标区域
                    icon_text = point.get("title", "")[:2]
                    icon_box = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(x + 1.3), Inches(1.8), Inches(1), Inches(1)
                    )
                    icon_box.fill.solid()
                    icon_box.fill.fore_color.rgb = COLORS['primary'] if i == 0 else (COLORS['secondary'] if i == 1 else COLORS['accent'])
                    icon_box.line.fill.background()
                    
                    # 标题
                    point_title = point.get("title", "")
                    title_box = slide.shapes.add_textbox(Inches(x), Inches(3.0), Inches(col_width), Inches(0.8))
                    title_frame = title_box.text_frame
                    title_frame.word_wrap = True
                    title_frame.text = point_title
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(18)
                    title_para.font.bold = True
                    title_para.font.color.rgb = COLORS['dark']
                    title_para.alignment = PP_ALIGN.CENTER
                    
                    # 详细说明
                    detail = point.get("detail", "")
                    detail_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.8), Inches(col_width - 0.4), Inches(2.5))
                    detail_frame = detail_box.text_frame
                    detail_frame.word_wrap = True
                    detail_frame.text = detail
                    detail_para = detail_frame.paragraphs[0]
                    detail_para.font.size = Pt(13)
                    detail_para.font.color.rgb = COLORS['dark']
                    detail_para.line_spacing = 1.3
                    
            elif layout_type == "left_text_right_chart":
                # 左文右图布局
                # 左侧要点
                for i, point in enumerate(main_points[:3]):
                    y = 1.8 + i * 1.5
                    
                    # 要点标记
                    bullet = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(0.7), Inches(y + 0.1), Inches(0.2), Inches(0.2)
                    )
                    bullet.fill.solid()
                    bullet.fill.fore_color.rgb = COLORS['primary']
                    bullet.line.fill.background()
                    
                    # 标题
                    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(5.5), Inches(0.5))
                    title_frame = title_box.text_frame
                    title_frame.text = point.get("title", "")
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(16)
                    title_para.font.bold = True
                    title_para.font.color.rgb = COLORS['dark']
                    
                    # 详细说明
                    detail_box = slide.shapes.add_textbox(Inches(1.1), Inches(y + 0.5), Inches(5.5), Inches(0.9))
                    detail_frame = detail_box.text_frame
                    detail_frame.word_wrap = True
                    detail_frame.text = point.get("detail", "")
                    detail_para = detail_frame.paragraphs[0]
                    detail_para.font.size = Pt(12)
                    detail_para.font.color.rgb = COLORS['dark']
                    detail_para.line_spacing = 1.2
                
                # 右侧图表区域
                chart_data = visual_suggestions.get("chart_data", {})
                if chart_data:
                    add_chart(slide, 7, 1.8, 5.5, 4, 'bar', chart_data, COLORS)
                else:
                    # 数字卡片
                    for i in range(min(2, len(main_points))):
                        highlight = main_points[i].get("highlight", "")
                        if highlight:
                            add_number_card(slide, 7.5 + i * 2.8, 2.5, 2.5, 1.5, 
                                          highlight, main_points[i].get("title", "")[:6], 
                                          COLORS['primary'], COLORS['light'])
                            
            else:
                # 标准布局 - 列表形式
                for i, point in enumerate(main_points[:5]):
                    y = 1.8 + i * 1.1
                    
                    # 序号圆圈
                    num_circle = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, Inches(0.6), Inches(y + 0.05), Inches(0.4), Inches(0.4)
                    )
                    num_circle.fill.solid()
                    num_circle.fill.fore_color.rgb = COLORS['primary']
                    num_circle.line.fill.background()
                    
                    # 序号文字
                    num_text = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.05), Inches(0.4), Inches(0.4))
                    num_text_frame = num_text.text_frame
                    num_text_frame.text = str(i + 1)
                    num_para = num_text_frame.paragraphs[0]
                    num_para.font.size = Pt(14)
                    num_para.font.bold = True
                    num_para.font.color.rgb = COLORS['white']
                    num_para.alignment = PP_ALIGN.CENTER
                    
                    # 标题
                    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11.5), Inches(0.5))
                    title_frame = title_box.text_frame
                    title_frame.text = point.get("title", "")
                    title_para = title_frame.paragraphs[0]
                    title_para.font.size = Pt(18)
                    title_para.font.bold = True
                    title_para.font.color.rgb = COLORS['dark']
                    
                    # 详细说明
                    detail_box = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.5), Inches(11.5), Inches(0.6))
                    detail_frame = detail_box.text_frame
                    detail_frame.word_wrap = True
                    detail_frame.text = point.get("detail", "")
                    detail_para = detail_frame.paragraphs[0]
                    detail_para.font.size = Pt(13)
                    detail_para.font.color.rgb = COLORS['dark']
                    detail_para.line_spacing = 1.2
                    
        elif slide_type == "summary":
            # 总结页 - 卡片式设计
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLORS['light']
            bg.line.fill.background()
            
            # 标题
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['dark']
            
            # 总结要点 - 2x2卡片布局
            for i, point in enumerate(main_points[:4]):
                col = i % 2
                row = i // 2
                x = 0.6 + col * 6.3
                y = 1.5 + row * 2.8
                
                # 卡片背景
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6), Inches(2.5)
                )
                card.fill.solid()
                card.fill.fore_color.rgb = COLORS['white']
                card.line.color.rgb = RGBColor(222, 226, 230)
                
                # 左侧色条
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.1), Inches(2.5)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = COLORS['primary'] if i % 2 == 0 else COLORS['secondary']
                bar.line.fill.background()
                
                # 勾选图标
                check = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.3), Inches(0.5), Inches(0.5)
                )
                check.fill.solid()
                check.fill.fore_color.rgb = COLORS['primary'] if i % 2 == 0 else COLORS['secondary']
                check.line.fill.background()
                
                # 标题
                card_title = slide.shapes.add_textbox(Inches(x + 0.9), Inches(y + 0.2), Inches(4.8), Inches(0.6))
                card_title_frame = card_title.text_frame
                card_title_frame.text = point.get("title", "")
                card_title_para = card_title_frame.paragraphs[0]
                card_title_para.font.size = Pt(16)
                card_title_para.font.bold = True
                card_title_para.font.color.rgb = COLORS['dark']
                
                # 详细内容
                card_detail = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.9), Inches(5.4), Inches(1.4))
                card_detail_frame = card_detail.text_frame
                card_detail_frame.word_wrap = True
                card_detail_frame.text = point.get("detail", "")
                card_detail_para = card_detail_frame.paragraphs[0]
                card_detail_para.font.size = Pt(12)
                card_detail_para.font.color.rgb = COLORS['dark']
                card_detail_para.line_spacing = 1.2
                    
        elif slide_type == "ending":
            # 结束页
            bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = COLORS['background']
            bg.line.fill.background()
            
            # 装饰线条
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.2), Inches(5.333), Inches(0.04)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['white']
            line.line.fill.background()
            
            # 主标题
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(1.2))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(64)
            title_para.font.bold = True
            title_para.font.color.rgb = COLORS['white']
            title_para.alignment = PP_ALIGN.CENTER
            
            # 副标题
            if subtitle:
                subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.733), Inches(0.8))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = subtitle
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(24)
                subtitle_para.font.color.rgb = COLORS['white']
                subtitle_para.alignment = PP_ALIGN.CENTER
            
            # 底部装饰
            bottom_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7)
            )
            bottom_bar.fill.solid()
            bottom_bar.fill.fore_color.rgb = COLORS['secondary']
            bottom_bar.line.fill.background()
    
    prs.save(output_path)
    return output_path


async def upload_ppt_to_feishu(file_path: str, title: str) -> str:
    """上传PPT文件到飞书云空间"""
    try:
        token = await feishu_api.get_tenant_access_token()
        upload_url = "https://open.feishu.cn/open-apis/drive/v1/files/upload_all"
        
        with open(file_path, 'rb') as f:
            file_content = f.read()
        
        file_size = len(file_content)
        file_name = f"{title}.pptx"
        
        import httpx
        
        headers = {"Authorization": f"Bearer {token}"}
        
        files = {
            'file': (file_name, file_content, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        }
        
        data = {
            'file_name': file_name,
            'parent_type': 'explorer',
            'parent_node': 'root',
            'size': str(file_size)
        }
        
        async with httpx.AsyncClient() as client:
            response = await client.post(upload_url, headers=headers, data=data, files=files)
            data = response.json()
            
            if data.get("code") == 0:
                file_token = data["data"]["file_token"]
                file_url = f"https://open.feishu.cn/open-apis/drive/v1/files/{file_token}"
                return file_url
            else:
                logger.error(f"上传PPT到飞书失败: {data.get('msg')}")
                return f"file://{file_path}"
                
    except Exception as e:
        logger.error(f"上传PPT到飞书异常: {e}")
        return f"file://{file_path}"


def format_ppt_satisfaction_check(ppt_url: str, ppt_title: str, revision_count: int = 0) -> str:
    """格式化PPT满意度确认界面"""
    lines = []
    lines.append("=" * 60)
    lines.append("✅ PPT生成完成！")
    lines.append("=" * 60)
    lines.append("")
    lines.append(f"📄 标题: {ppt_title}")
    lines.append(f"🔗 链接: {ppt_url}")
    if revision_count > 0:
        lines.append(f"📝 已修改次数: {revision_count}")
    lines.append("")
    lines.append("请查看PPT后选择：")
    lines.append("")
    lines.append("  [1] 满意 - PPT符合要求，完成任务")
    lines.append("  [2] 需要修改 - 在现有PPT基础上调整")
    lines.append("  [3] 重新生成 - 完全重做（保留大纲，重新生成内容）")
    lines.append("")
    lines.append("=" * 60)
    
    return "\n".join(lines)


async def revise_ppt_content(state: IMState) -> IMState:
    """根据用户反馈修改PPT内容（在原有基础上修改）"""
    ppt_content = state.get("ppt_content")
    ppt_structure = state.get("ppt_structure")
    feedback = state.get("ppt_satisfaction_feedback", "")
    chat_context = state.get("chat_context", "")
    
    if not ppt_content or not feedback:
        logger.error("[ppt_generate_node] 缺少PPT内容或修改意见")
        state["error"] = "缺少PPT内容或修改意见"
        return state
    
    try:
        logger.info(f"[ppt_generate_node] 根据用户反馈修改PPT: {feedback}")
        state["messages"].append(f"[ppt_generate_node] 正在根据反馈修改PPT...")
        
        # 构建修改消息
        revision_prompt = f"""你是一个专业的PPT内容修改专家。用户已经对生成的PPT提出了修改意见，请根据意见修改PPT内容。

用户修改意见：{feedback}

请分析用户的修改意见，并逐页调整PPT内容。修改原则：
1. 保留原有结构和风格
2. 根据用户意见增删或调整内容
3. 保持专业性和可读性
4. 确保修改后的内容更加符合用户需求

请输出修改后的完整PPT内容JSON。"""

        # 将当前PPT内容转换为JSON字符串
        current_content_json = json.dumps(ppt_content, ensure_ascii=False, indent=2)
        
        messages = [
            SystemMessage(content=revision_prompt),
            HumanMessage(content=f"当前PPT内容：\n{current_content_json}\n\n请根据用户意见修改。")
        ]
        
        # 调用LLM进行修改
        res = await ppt_content_agent.ainvoke({"messages": messages})
        raw_content = res["messages"][-1].content
        json_str = extract_json(raw_content)
        
        try:
            revised_content = json.loads(json_str)
            state["ppt_content"] = revised_content
            logger.info("[ppt_generate_node] PPT内容修改完成")
            state["messages"].append("[ppt_generate_node] PPT内容修改完成")
        except json.JSONDecodeError as e:
            logger.warning(f"[ppt_generate_node] 修改后内容解析失败，使用原内容: {e}")
            state["messages"].append("[ppt_generate_node] 修改解析失败，将重新渲染原内容")
        
        # 重新渲染PPT文件
        logger.info("[ppt_generate_node] 重新渲染修改后的PPT")
        state["messages"].append("[ppt_generate_node] 重新渲染PPT文件...")
        
        output_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(output_dir, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w\s-]', '', ppt_structure.get("title", "未命名PPT")).strip()
        revision_count = state.get("ppt_revision_count", 0)
        file_name = f"{safe_title}_修订{revision_count}_{timestamp}.pptx"
        file_path = os.path.join(output_dir, file_name)
        
        create_ppt_file(ppt_structure, state["ppt_content"], file_path)
        logger.info(f"[ppt_generate_node] 修改后的PPT文件已生成: {file_path}")
        state["messages"].append(f"[ppt_generate_node] 修改后的PPT文件已生成")
        
        # 重新上传
        ppt_url = await upload_ppt_to_feishu(file_path, ppt_structure.get("title", "未命名PPT"))
        state["ppt_url"] = ppt_url
        
        # 增加修改次数
        state["ppt_revision_count"] = revision_count + 1
        
        logger.info(f"[ppt_generate_node] 修改后的PPT上传完成: {ppt_url}")
        state["messages"].append(f"[ppt_generate_node] 修改后的PPT链接: {ppt_url}")
        
        # 重置满意度确认状态，准备下一次确认
        state["ppt_satisfaction_confirmed"] = False
        state["confirmed"] = False
        state["need_confirm"] = True
        state["confirm_type"] = "ppt_satisfaction"
        
    except Exception as e:
        logger.error(f"[ppt_generate_node] PPT修改异常: {e}")
        state["error"] = f"PPT修改失败: {str(e)}"
        state["messages"].append(f"[ppt_generate_node] PPT修改失败: {str(e)}")
    
    return state


async def generate_ppt_content(state: IMState) -> IMState:
    """生成PPT详细内容并渲染文件（支持满意度确认）"""
    ppt_structure = state.get("ppt_structure")
    chat_context = state.get("chat_context", "")
    
    if not ppt_structure:
        logger.error("[ppt_generate_node] 缺少PPT大纲，无法生成内容")
        state["error"] = "缺少PPT大纲，无法生成内容"
        return state
    
    # 检查是否是修改后的重新生成
    satisfaction_feedback = state.get("ppt_satisfaction_feedback")
    if satisfaction_feedback and state.get("ppt_content"):
        # 用户要求重新生成，但保留大纲
        logger.info("[ppt_generate_node] 用户要求重新生成PPT内容（保留大纲）")
        state["messages"].append("[ppt_generate_node] 根据反馈重新生成PPT内容...")
        # 清除旧内容，但保留大纲
        state["ppt_content"] = None
    
    # 检查是否已生成过内容（用于满意度确认流程）
    if state.get("ppt_content") and state.get("ppt_url"):
        # 已经生成过，进入满意度确认流程
        if not state.get("ppt_satisfaction_confirmed", False):
            logger.info("[ppt_generate_node] PPT已生成，等待用户满意度确认")
            state["messages"].append("[ppt_generate_node] PPT已生成，等待用户确认满意度")
            
            ppt_url = state.get("ppt_url", "")
            ppt_title = ppt_structure.get("title", "未命名PPT")
            revision_count = state.get("ppt_revision_count", 0)
            
            satisfaction_menu = format_ppt_satisfaction_check(ppt_url, ppt_title, revision_count)
            state["messages"].append(satisfaction_menu)
            
            # 设置确认状态
            state["current_scene_before_confirm"] = "ppt_generate_node"
            state["need_confirm"] = True
            state["confirm_type"] = "ppt_satisfaction"
            state["confirmed"] = False
            
            return state
        else:
            # 用户已确认满意，流程结束
            logger.info("[ppt_generate_node] 用户已确认PPT满意")
            state["messages"].append("[ppt_generate_node] 用户确认PPT满意，任务完成")
            return state
    
    try:
        logger.info("[ppt_generate_node] Step D2: 逐页生成PPT内容")
        state["messages"].append("[ppt_generate_node] 逐页生成PPT内容...")
        
        slides = ppt_structure.get("slides", [])
        refined_slides = []
        
        for slide in slides:
            messages = build_ppt_content_messages(ppt_structure, slide, chat_context)
            res = await ppt_content_agent.ainvoke({"messages": messages})
            
            raw_content = res["messages"][-1].content
            json_str = extract_json(raw_content)
            
            try:
                slide_content = json.loads(json_str)
                refined_slides.append(slide_content)
            except json.JSONDecodeError as e:
                logger.warning(f"[ppt_generate_node] 页面 {slide.get('page')} 内容解析失败，使用原始内容")
                # 转换旧格式到新格式
                old_content = slide.get("content", {})
                points = old_content.get("points", [])
                refined_slides.append({
                    "page": slide.get("page"),
                    "type": slide.get("type"),
                    "title": slide.get("title"),
                    "subtitle": slide.get("subtitle", ""),
                    "content": {
                        "main_points": [{"title": p, "detail": "", "highlight": ""} for p in points] if points else [],
                        "visual_suggestions": {"layout": "standard"},
                        "notes": ""
                    }
                })
        
        ppt_content = {
            "title": ppt_structure.get("title"),
            "ppt_type": ppt_structure.get("ppt_type"),
            "ppt_style": ppt_structure.get("ppt_style", "business_blue"),
            "slides": refined_slides
        }
        state["ppt_content"] = ppt_content
        logger.info(f"[ppt_generate_node] PPT内容生成完成，共 {len(refined_slides)} 页")
        state["messages"].append(f"[ppt_generate_node] PPT内容生成完成，共 {len(refined_slides)} 页")
        
        logger.info("[ppt_generate_node] Step D3: 渲染PPT文件")
        state["messages"].append("[ppt_generate_node] 渲染PPT文件...")
        
        output_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(output_dir, exist_ok=True)
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = re.sub(r'[^\w\s-]', '', ppt_structure.get("title", "未命名PPT")).strip()
        file_name = f"{safe_title}_{timestamp}.pptx"
        file_path = os.path.join(output_dir, file_name)
        
        create_ppt_file(ppt_structure, ppt_content, file_path)
        logger.info(f"[ppt_generate_node] PPT文件已生成: {file_path}")
        state["messages"].append(f"[ppt_generate_node] PPT文件已生成: {file_path}")
        
        logger.info("[ppt_generate_node] Step D4: 上传PPT到飞书")
        state["messages"].append("[ppt_generate_node] 上传PPT到飞书...")
        
        ppt_url = await upload_ppt_to_feishu(file_path, ppt_structure.get("title", "未命名PPT"))
        state["ppt_url"] = ppt_url
        
        logger.info(f"[ppt_generate_node] PPT上传完成: {ppt_url}")
        state["messages"].append(f"[ppt_generate_node] PPT生成完成，链接: {ppt_url}")
        
        # 进入满意度确认流程
        state["ppt_satisfaction_confirmed"] = False
        state["confirmed"] = False
        state["need_confirm"] = True
        state["confirm_type"] = "ppt_satisfaction"
        state["ppt_revision_count"] = 0  # 初始化修改次数
        
        # 显示满意度确认界面
        satisfaction_menu = format_ppt_satisfaction_check(ppt_url, ppt_structure.get("title", "未命名PPT"), 0)
        state["messages"].append(satisfaction_menu)
        
    except Exception as e:
        logger.error(f"[ppt_generate_node] PPT内容生成异常: {e}")
        state["error"] = f"PPT内容生成失败: {str(e)}"
        state["messages"].append(f"[ppt_generate_node] PPT内容生成失败: {str(e)}")
    
    return state


# PPT风格选项定义
PPT_STYLE_OPTIONS = {
    "business_blue": {
        "name": "商务蓝",
        "description": "专业稳重，适合商务汇报、企业介绍",
        "emoji": "💼",
        "features": ["蓝色主调", "正式专业", "数据驱动"]
    },
    "tech_dark": {
        "name": "科技黑",
        "description": "科技感强，适合技术分享、产品发布",
        "emoji": "💻",
        "features": ["深色背景", "科技元素", "代码展示"]
    },
    "fresh_green": {
        "name": "清新绿",
        "description": "亲和自然，适合环保、教育、健康主题",
        "emoji": "🌿",
        "features": ["绿色主调", "清新自然", "亲和力强"]
    },
    "warm_orange": {
        "name": "活力橙",
        "description": "活力热情，适合创意、营销、活动推广",
        "emoji": "🔥",
        "features": ["橙色主调", "活力创意", "吸引眼球"]
    },
    "minimal_white": {
        "name": "极简白",
        "description": "简洁干净，适合学术、研究、极简风格",
        "emoji": "📄",
        "features": ["白色为主", "极简设计", "内容聚焦"]
    },
    "elegant_purple": {
        "name": "优雅紫",
        "description": "高端精致，适合时尚、艺术、品牌展示",
        "emoji": "💜",
        "features": ["紫色渐变", "优雅精致", "品质感强"]
    }
}


def format_style_selection() -> str:
    """格式化风格选择界面"""
    lines = []
    lines.append("=" * 60)
    lines.append("🎨 请选择PPT风格模板")
    lines.append("=" * 60)
    lines.append("")
    
    for key, style in PPT_STYLE_OPTIONS.items():
        lines.append(f"{style['emoji']} [{key}] {style['name']}")
        lines.append(f"   描述: {style['description']}")
        lines.append(f"   特点: {', '.join(style['features'])}")
        lines.append("")
    
    lines.append("=" * 60)
    lines.append("请输入风格名称（如: business_blue）或序号（1-6）")
    lines.append("• 输入 'auto' 让系统自动推荐")
    lines.append("• 输入 'cancel' 取消PPT生成")
    lines.append("=" * 60)
    
    return "\n".join(lines)


def parse_style_selection(user_input: str) -> tuple:
    """解析用户风格选择
    
    Returns:
        (style_key, is_valid, message)
    """
    user_input = user_input.strip().lower()
    
    if user_input in ['cancel', '取消', 'q', 'quit']:
        return None, False, "cancelled"
    
    if user_input in ['auto', '自动', '推荐']:
        return "auto", True, "将自动根据主题推荐风格"
    
    # 检查是否是数字序号
    try:
        num = int(user_input)
        style_keys = list(PPT_STYLE_OPTIONS.keys())
        if 1 <= num <= len(style_keys):
            return style_keys[num - 1], True, f"已选择: {PPT_STYLE_OPTIONS[style_keys[num - 1]]['name']}"
    except ValueError:
        pass
    
    # 检查是否是风格名称
    if user_input in PPT_STYLE_OPTIONS:
        return user_input, True, f"已选择: {PPT_STYLE_OPTIONS[user_input]['name']}"
    
    # 模糊匹配
    for key, style in PPT_STYLE_OPTIONS.items():
        if user_input in style['name'].lower() or user_input in key.lower():
            return key, True, f"已选择: {style['name']}"
    
    return None, False, f"无效的选择 '{user_input}'，请重新输入"


async def ppt_generate_node(state: IMState) -> IMState:
    """场景D：PPT生成
    
    完整执行流程:
    1. 检查用户是否已选择风格
       - 未选择: 显示风格选项，等待用户选择
       - 已选择: 继续下一步
    2. 检查用户是否已确认大纲
       - 未确认: 生成大纲，等待用户确认
       - 已确认: 执行内容生成
    3. 检查用户满意度
       - 未确认: 显示满意度确认界面
       - 满意: 任务完成
       - 需要修改: 在原有基础上修改
       - 重新生成: 保留大纲，重新生成内容
    """
    state["current_scene"] = "ppt_generate_node"
    state["messages"].append("[ppt_generate_node] 进入PPT生成节点")
    
    if state.get("cancelled"):
        logger.info("[ppt_generate_node] 用户已取消任务，跳过PPT生成")
        state["messages"].append("[ppt_generate_node] 用户取消任务，跳过执行")
        return state
    
    # Step 1: 检查风格选择
    ppt_style_selected = state.get("ppt_style_selected")
    ppt_style_confirmed = state.get("ppt_style_confirmed", False)
    
    if not ppt_style_confirmed:
        # 用户还未选择风格，显示风格选项
        logger.info("[ppt_generate_node] 等待用户选择PPT风格")
        state["messages"].append("[ppt_generate_node] 请用户选择PPT风格")
        
        style_menu = format_style_selection()
        state["messages"].append(style_menu)
        
        # 设置状态，让ConfirmNode知道这是风格选择
        state["current_scene_before_confirm"] = "ppt_generate_node"
        state["need_confirm"] = True
        state["confirm_type"] = "style_selection"
        state["confirmed"] = False
        
        return state
    
    # Step 2: 检查是否需要修改PPT（满意度确认后的修改）
    revision_type = state.get("ppt_revision_type")
    satisfaction_feedback = state.get("ppt_satisfaction_feedback")
    
    if revision_type == "revise" and satisfaction_feedback:
        # 用户要求在原有基础上修改
        logger.info("[ppt_generate_node] 用户要求修改PPT内容")
        state["messages"].append("[ppt_generate_node] 开始修改PPT内容")
        
        # 清除修改类型标记，避免重复修改
        state["ppt_revision_type"] = None
        
        return await revise_ppt_content(state)
    
    if revision_type == "regenerate" and satisfaction_feedback:
        # 用户要求重新生成（保留大纲）
        logger.info("[ppt_generate_node] 用户要求重新生成PPT")
        state["messages"].append("[ppt_generate_node] 开始重新生成PPT内容")
        
        # 清除修改类型标记和内容
        state["ppt_revision_type"] = None
        state["ppt_content"] = None
        state["ppt_url"] = ""
        
        # 重新生成内容
        return await generate_ppt_content(state)
    
    # Step 3: 检查大纲确认
    ppt_structure = state.get("ppt_structure")
    confirmed = state.get("confirmed", False)
    ppt_outline_feedback = state.get("ppt_outline_feedback")

    if ppt_structure and confirmed:
        # 用户已确认大纲，开始生成内容
        # 但先检查是否已经生成过内容（满意度确认流程）
        if state.get("ppt_content") and state.get("ppt_url"):
            # 已经生成过内容，进入满意度确认
            return await generate_ppt_content(state)
        
        logger.info("[ppt_generate_node] 用户已确认大纲，开始生成PPT内容")
        state["messages"].append("[ppt_generate_node] 用户确认大纲，开始生成内容")
        return await generate_ppt_content(state)
    
    if ppt_structure and ppt_outline_feedback:
        logger.info("[ppt_generate_node] 用户要求修改大纲，重新生成")
        state["messages"].append("[ppt_generate_node] 根据用户反馈重新生成大纲")
        state["confirmed"] = False
        state["need_confirm"] = True
        state["confirm_type"] = "outline_confirmation"
    else:
        logger.info("[ppt_generate_node] 首次生成PPT大纲")
        state["messages"].append("[ppt_generate_node] 首次生成PPT大纲")
    
    try:
        messages = build_ppt_outline_messages(state)
        
        if ppt_outline_feedback and ppt_structure:
            res = await ppt_outline_revision_agent.ainvoke({"messages": messages})
        else:
            res = await ppt_outline_agent.ainvoke({"messages": messages})
        
        raw_content = res["messages"][-1].content
        json_str = extract_json(raw_content)
        
        try:
            ppt_structure = json.loads(json_str)
        except json.JSONDecodeError as e:
            logger.error(f"[ppt_generate_node] PPT结构JSON解析失败: {e}")
            logger.error(f"[ppt_generate_node] 原始内容: {raw_content}")
            state["error"] = f"PPT结构解析失败: {e}"
            return state
        
        state["ppt_structure"] = ppt_structure
        state["current_scene_before_confirm"] = "ppt_generate_node"
        state["need_confirm"] = True
        state["confirmed"] = False
        state["confirm_type"] = "outline_confirmation"
        
        state.pop("ppt_outline_feedback", None)
        
        # 显示风格信息
        selected_style = state.get("ppt_style_selected", "business_blue")
        style_name = PPT_STYLE_OPTIONS.get(selected_style, {}).get("name", "商务蓝")
        
        formatted_outline = format_ppt_outline(ppt_structure)
        logger.info(f"[ppt_generate_node] PPT大纲生成完成 (风格: {style_name}):\n{formatted_outline}")
        state["messages"].append(f"[ppt_generate_node] 已选择风格: {style_name}")
        state["messages"].append(f"[ppt_generate_node] PPT大纲生成完成，等待用户确认")
        state["messages"].append(formatted_outline)
        
    except Exception as e:
        logger.error(f"[ppt_generate_node] PPT大纲生成异常: {e}")
        state["error"] = f"PPT大纲生成失败: {str(e)}"
        state["messages"].append(f"[ppt_generate_node] PPT大纲生成失败: {str(e)}")
    
    return state
